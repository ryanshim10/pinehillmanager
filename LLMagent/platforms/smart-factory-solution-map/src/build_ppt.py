"""Build Smart Factory Solution Map PPT.

Policy:
- Only write verifiable facts.
- Every non-trivial statement must have a citation.
- If missing citation: leave blank.

This script generates a PPT skeleton + placeholders for citations.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# Style (Bain-like)
COLOR_RED = RGBColor(192, 0, 0)
COLOR_GREY = RGBColor(110, 110, 110)
COLOR_DARK = RGBColor(30, 30, 30)
FONT_FAMILY = "Malgun Gothic"

# Layout constants
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.35)
MARGIN_B = Inches(0.35)
LEFT_COL_W = Inches(2.2)
TOP_META_H = Inches(0.35)


@dataclass
class Citation:
    label: str
    url: str
    accessed_utc: str


@dataclass
class SlideSection:
    title: str
    bullets: List[str]
    citations: List[Citation]


GROUPS = [
    "SK",
    "삼성",
    "LG",
    "포스코",
    "현대자동차그룹",
    "HD현대",
    "두산그룹",
    "LS그룹",
]


def _set_text_run(run, *, size=14, bold=False, color=COLOR_DARK, name=FONT_FAMILY):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_meta_header(slide, meta: str, page_no: int | None = None):
    # Top meta line
    box = slide.shapes.add_textbox(MARGIN_L, MARGIN_T, Inches(9.0), TOP_META_H)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = meta
    _set_text_run(r, size=10, color=COLOR_GREY)

    # Thin red line under header
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_L,
        MARGIN_T + TOP_META_H + Inches(0.05),
        Inches(9.0),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_RED
    line.line.color.rgb = COLOR_RED

    # Footer page number
    if page_no is not None:
        f = slide.shapes.add_textbox(Inches(9.2), Inches(7.05), Inches(0.5), Inches(0.3))
        tf2 = f.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = str(page_no)
        _set_text_run(r2, size=10, color=COLOR_GREY)


def add_bain_bullets_slide(
    prs: Presentation,
    *,
    title: str,
    left_label: str,
    bullets: List[str],
    meta: str = "Source: Official sites / Press releases (verify per bullet)",
    sources: List[str] | None = None,
    page_no: int | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _add_meta_header(slide, meta, page_no=page_no)

    # Left icon circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        MARGIN_L,
        Inches(1.7),
        Inches(0.75),
        Inches(0.75),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = COLOR_RED
    circ.line.color.rgb = COLOR_RED

    # Left label
    lb = slide.shapes.add_textbox(MARGIN_L + Inches(0.9), Inches(1.65), LEFT_COL_W - Inches(0.9), Inches(1.0))
    tf = lb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = left_label
    _set_text_run(r, size=18, bold=True, color=COLOR_RED)

    # Vertical red bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_L + LEFT_COL_W,
        Inches(1.55),
        Inches(0.08),
        Inches(3.7),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_RED
    bar.line.color.rgb = COLOR_RED

    # Title
    tb_title = slide.shapes.add_textbox(MARGIN_L, Inches(1.05), Inches(9.0), Inches(0.6))
    tf_t = tb_title.text_frame
    tf_t.clear()
    p_t = tf_t.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = title
    _set_text_run(r_t, size=22, bold=True, color=COLOR_DARK)

    # Bullets
    content_x = MARGIN_L + LEFT_COL_W + Inches(0.25)
    content_w = Inches(9.0) - LEFT_COL_W - Inches(0.25)
    content = slide.shapes.add_textbox(content_x, Inches(1.75), content_w, Inches(4.0))
    tfc = content.text_frame
    tfc.clear()
    for i, b in enumerate(bullets):
        if not b.strip():
            continue
        p = tfc.paragraphs[0] if i == 0 else tfc.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = b
        _set_text_run(r, size=16, color=COLOR_DARK)

    # Sources small box
    sources = sources or []
    src_box = slide.shapes.add_textbox(content_x, Inches(6.0), content_w, Inches(1.0))
    tfs = src_box.text_frame
    tfs.clear()
    p0 = tfs.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "Sources:"
    _set_text_run(r0, size=10, bold=True, color=COLOR_GREY)
    for s in sources[:3]:
        p = tfs.add_paragraph()
        r = p.add_run()
        r.text = f"- {s}"
        _set_text_run(r, size=9, color=COLOR_GREY)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_meta_header(slide, "Source: Lit. Search, Official Sites, Press Releases", page_no=1)

    tb = slide.shapes.add_textbox(MARGIN_L, Inches(2.2), Inches(9.0), Inches(1.5))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_text_run(r, size=34, bold=True, color=COLOR_DARK)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    _set_text_run(r2, size=16, color=COLOR_GREY)


def build() -> Path:
    prs = Presentation()

    add_title_slide(
        prs,
        title="Smart Factory Solution Map (Draft)",
        subtitle="대상: SK / 삼성 / LG / 포스코 / 현대자동차그룹 / HD현대 / 두산그룹 / LS그룹\n원칙: 출처 없는 내용은 작성하지 않음",
    )

    add_bain_bullets_slide(
        prs,
        title="작성 원칙",
        left_label="원칙",
        bullets=[
            "팩트만 작성 (출처 필수)",
            "불확실/상충 정보는 제외",
            "[추정]은 근거 출처가 있을 때만",
            "PPT에는 요약, MD에는 상세 기록",
        ],
        sources=["LLMagent/platforms/smart-factory-solution-map/README.md (2026-02-06 UTC)"],
        page_no=2,
    )

    page_no = 3
    for g in GROUPS:
        add_bain_bullets_slide(
            prs,
            title=f"{g} | 개요",
            left_label=g,
            bullets=[
                "(요약 작성 중) 그룹별 스마트팩토리/제조DX 방향",
                "(요약 작성 중) 대표 계열사/플랫폼/거점",
            ],
            sources=[f"platforms/smart-factory-solution-map/sources/{g}.md"],
            page_no=page_no,
        )
        page_no += 1

        add_bain_bullets_slide(
            prs,
            title=f"{g} | 외부 생태계 맵 A (파트너 카테고리)",
            left_label="A",
            bullets=[
                "클라우드/인프라 파트너:",
                "산업 SW (MES/PLM/APS) 파트너:",
                "데이터/AI 파트너:",
                "로봇/자동화 파트너:",
                "머신비전/검사 파트너:",
                "OT/보안 파트너:",
                "기타(통신/5G/Edge 등):",
            ],
            sources=[],
            page_no=page_no,
        )
        page_no += 1

        add_bain_bullets_slide(
            prs,
            title=f"{g} | 외부 생태계 맵 B (밸류체인)",
            left_label="B",
            bullets=[
                "데이터 수집(센서/PLC/SCADA):",
                "데이터 파이프라인/저장:",
                "분석/AI(MLOps 포함):",
                "실행/제어(현장 적용):",
                "업무/운영(MES/품질/설비/에너지):",
                "가상화(디지털트윈/시뮬):",
                "보안/거버넌스:",
            ],
            sources=[],
            page_no=page_no,
        )
        page_no += 1

        add_bain_bullets_slide(
            prs,
            title=f"{g} | 솔루션/역량 맵 (요약)",
            left_label="맵",
            bullets=[
                "MES/PLM",
                "IIoT/설비데이터",
                "AI/분석",
                "로봇/자동화",
                "머신비전/검사",
                "디지털트윈/시뮬레이션",
                "보안/OT",
            ],
            sources=[],
            page_no=page_no,
        )
        page_no += 1

        # Sources slide (placeholder; will be filled from sources/*.md later)
        add_bain_bullets_slide(
            prs,
            title=f"{g} | Sources",
            left_label="출처",
            bullets=[
                f"(작성 중) platforms/smart-factory-solution-map/sources/{g}.md 참고",
            ],
            sources=[f"platforms/smart-factory-solution-map/sources/{g}.md"],
            page_no=page_no,
        )
        page_no += 1

    out = OUT_DIR / "SmartFactory_SolutionMap_8groups.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(str(path))

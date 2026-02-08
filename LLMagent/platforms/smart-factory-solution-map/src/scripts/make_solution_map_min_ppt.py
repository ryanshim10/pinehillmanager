#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

PRIMARY = rgb('#1E5EFF')
PRIMARY_DARK = rgb('#0B1B4D')
LIGHT = rgb('#EAF2FF')
TEXT = rgb('#111827')
WHITE = rgb('#FFFFFF')
MUTED = rgb('#4B5563')


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    return shp


def add_topbar(slide, prs, title):
    rect(slide, 0, 0, prs.slide_width, Inches(0.7), PRIMARY)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), prs.slide_width-Inches(1.2), Inches(0.5))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE


def add_title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PRIMARY_DARK)
    rect(s, 0, Inches(5.8), prs.slide_width, Inches(1.7), PRIMARY)

    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.8), prs.slide_width-Inches(1.8), Inches(1.8))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(44); r.font.bold=True; r.font.color.rgb = WHITE

    sb = s.shapes.add_textbox(Inches(0.9), Inches(3.6), prs.slide_width-Inches(1.8), Inches(1.2))
    stf = sb.text_frame; stf.clear()
    p2 = stf.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT


def add_group_slide(prs, group, a_items, b_items, sources):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_topbar(s, prs, f"{group} — 최소 완성본 (A/B 맵)")

    # Two columns
    margin = Inches(0.7)
    col_w = (prs.slide_width - margin*2 - Inches(0.4)) / 2
    col_h = prs.slide_height - Inches(1.2) - Inches(0.3)
    top = Inches(1.0)

    # A box
    rect(s, margin, top, col_w, col_h, LIGHT, line_rgb=PRIMARY)
    at = s.shapes.add_textbox(margin+Inches(0.3), top+Inches(0.2), col_w-Inches(0.6), Inches(0.4))
    tf = at.text_frame; tf.clear(); p=tf.paragraphs[0]
    p.text = "A. 파트너 카테고리(출처 있는 것만)"; p.font.size=Pt(18); p.font.bold=True; p.font.color.rgb=PRIMARY_DARK

    ab = s.shapes.add_textbox(margin+Inches(0.3), top+Inches(0.7), col_w-Inches(0.6), col_h-Inches(1.0))
    tf = ab.text_frame; tf.clear(); tf.word_wrap=True
    if not a_items:
        a_items = ["(출처 확보 중)"]
    for i, item in enumerate(a_items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT

    # B box
    left2 = margin + col_w + Inches(0.4)
    rect(s, left2, top, col_w, col_h, LIGHT, line_rgb=PRIMARY)
    bt = s.shapes.add_textbox(left2+Inches(0.3), top+Inches(0.2), col_w-Inches(0.6), Inches(0.4))
    tf = bt.text_frame; tf.clear(); p=tf.paragraphs[0]
    p.text = "B. 밸류체인(출처 있는 것만)"; p.font.size=Pt(18); p.font.bold=True; p.font.color.rgb=PRIMARY_DARK

    bb = s.shapes.add_textbox(left2+Inches(0.3), top+Inches(0.7), col_w-Inches(0.6), col_h-Inches(1.0))
    tf = bb.text_frame; tf.clear(); tf.word_wrap=True
    if not b_items:
        b_items = ["(출처 확보 중)"]
    for i, item in enumerate(b_items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT

    # Footer sources
    ft = s.shapes.add_textbox(margin, prs.slide_height-Inches(0.45), prs.slide_width-margin*2, Inches(0.35))
    tf = ft.text_frame; tf.clear(); p=tf.paragraphs[0]
    p.text = "Sources: " + (" | ".join(sources) if sources else "-")
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def _extract_section_lines(md: str, start_hdr: str, end_hdr: str):
    s = md.find(start_hdr)
    if s == -1:
        return []
    e = md.find(end_hdr, s + len(start_hdr))
    sec = md[s:e] if e != -1 else md[s:]
    out = []
    for ln in sec.splitlines():
        ln = ln.strip()
        # keep only top-level bullets
        if ln.startswith('- ') and 'Source:' not in ln and 'Fact:' not in ln and '추정' not in ln and 'To-' not in ln:
            out.append(ln[2:].strip())
    return out


def _extract_sources(md: str, limit: int = 5):
    src = []
    for ln in md.splitlines():
        ln = ln.strip()
        if ln.startswith('- Source:'):
            url = ln.split(':', 1)[1].strip()
            # shorten
            url = url.replace('https://', '').replace('http://', '')
            if url.endswith('/'):
                url = url[:-1]
            src.append(url)
    # de-dup keep order
    seen = set(); out = []
    for u in src:
        if u in seen:
            continue
        seen.add(u)
        out.append(u)
        if len(out) >= limit:
            break
    return out


def _extract_map_items(md: str, map_hdr: str, next_hdr: str):
    # In our sources logs, A/B sections contain mixed Fact/추정 lines.
    # For PPT we only include FACT-backed items: lines that are NOT marked '추정' and have a Source line nearby.
    s = md.find(map_hdr)
    if s == -1:
        return []
    e = md.find(next_hdr, s + len(map_hdr))
    sec = md[s:e] if e != -1 else md[s:]

    items = []
    cur = None
    cur_is_infer = False
    cur_has_source = False

    def flush():
        nonlocal cur, cur_is_infer, cur_has_source
        if cur and (not cur_is_infer) and cur_has_source:
            items.append(cur)
        cur = None
        cur_is_infer = False
        cur_has_source = False

    for ln in sec.splitlines():
        t = ln.strip()
        if t.startswith('- '):
            flush()
            txt = t[2:].strip()
            cur_is_infer = ('추정' in txt) or ('To-' in txt)
            # normalize prefixes
            for pref in ('Fact:', 'Fact', '팩트:', '팩트'):
                if txt.startswith(pref):
                    txt = txt.split(':', 1)[1].strip() if ':' in txt else txt[len(pref):].strip()
            # keep only short-ish summary
            cur = txt
        elif t.startswith('Source:') or t.startswith('- Source:'):
            cur_has_source = True

    flush()
    return items


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "스마트팩토리 솔루션 맵 — 최소 완성본", "8개 그룹(HD현대/두산/현대차그룹/삼성/LG/SK/포스코/LS그룹) | 출처 기반 요약")

    base = Path(__file__).resolve().parents[2] / 'sources'
    files = [
        ("HD현대", "HD현대.md"),
        ("두산", "두산그룹.md"),
        ("현대차그룹", "현대자동차그룹.md"),
        ("삼성", "삼성.md"),
        ("LG", "LG.md"),
        ("SK", "SK.md"),
        ("포스코", "포스코.md"),
        ("LS그룹", "LS그룹.md"),
    ]

    for group, fn in files:
        md = (base / fn).read_text(encoding='utf-8', errors='ignore')
        a_items = _extract_map_items(md, '## 3) 생태계 맵 A', '## 4) 생태계 맵 B')
        b_items = _extract_map_items(md, '## 4) 생태계 맵 B', '## 5) 미확인 / To-Verify')
        sources = _extract_sources(md, limit=6)
        add_group_slide(prs, group, a_items, b_items, sources)

    out_dir = Path(__file__).resolve().parents[2] / 'output'
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / 'smart-factory-solution-map_MIN_blue.pptx'
    prs.save(out_path)
    print(out_path)


if __name__ == '__main__':
    main()

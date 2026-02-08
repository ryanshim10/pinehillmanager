#!/usr/bin/env python3
"""Generate a denser Smart Factory Solution Map PPT.

Design goals:
- Use many small cards (category matrix) like vendor maps.
- Each card shows a short summary and links to a detailed slide.
- Separate Facts vs Inference (inference shown but labeled).
- KPI is always 4 lanes: P (Productivity/Leadtime), Q (Quality), A (Asset/OEE), E (Energy/ESG).

Slide structure per group:
- 1) Matrix (16 cards; each card hyperlinks to its detail slide if content exists)
- 2) KPI snapshot (P/Q/A/E)
- 3) Detail slides (only for categories with content)

Notes:
- Internal slide links are implemented via shape.click_action.target_slide.
- External links in the sources footer are clickable.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


PRIMARY = rgb('#1E5EFF')
PRIMARY_DARK = rgb('#0B1B4D')
LIGHT = rgb('#EAF2FF')
TEXT = rgb('#111827')
WHITE = rgb('#FFFFFF')
MUTED = rgb('#4B5563')
WARN = rgb('#B45309')


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    return shp


def add_topbar(slide, prs, title):
    rect(slide, 0, 0, prs.slide_width, Inches(0.65), PRIMARY)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), prs.slide_width - Inches(1.2), Inches(0.45))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PRIMARY_DARK)
    rect(s, 0, Inches(5.8), prs.slide_width, Inches(1.7), PRIMARY)

    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.8), prs.slide_width - Inches(1.8), Inches(1.6))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE

    sb = s.shapes.add_textbox(Inches(0.9), Inches(3.5), prs.slide_width - Inches(1.8), Inches(1.2))
    stf = sb.text_frame
    stf.clear()
    p2 = stf.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT


@dataclass
class FactItem:
    text: str
    kind: str  # 'fact' | 'infer'


CATEGORIES = [
    "MES/PLM",
    "HMI/SCADA",
    "IIoT/RTLS/Sensing",
    "Connectivity(5G/Net)",
    "Cloud/Data",
    "Data/AI",
    "Digital Twin/Sim",
    "Robot/AMR",
    "Vision/Inspection",
    "Energy/FEMS/ESG",
    "Safety",
    "OT Security",
    "Test/Metrology",
    "Ops/Workflow",
    "Integration(ERP)",
    "KPI/Outcome",
]


def categorize(text: str) -> str:
    t = text.lower()
    if any(k in t for k in ["mes", "plm"]):
        return "MES/PLM"
    if any(k in t for k in ["scada", "hmi"]):
        return "HMI/SCADA"
    if any(k in t for k in ["rtls", "위치", "센서", "iot", "수집"]):
        return "IIoT/RTLS/Sensing"
    if any(k in t for k in ["5g", "network", "네트워크", "통신", "kt"]):
        return "Connectivity(5G/Net)"
    if any(k in t for k in ["cloud", "클라우드", "data", "데이터"]):
        return "Cloud/Data"
    if any(k in t for k in ["ai", "인공지능", "예측", "분석"]):
        return "Data/AI"
    if any(k in t for k in ["digital twin", "디지털 트윈", "트윈포스", "3d모델", "시뮬레이션", "버추얼"]):
        return "Digital Twin/Sim"
    if any(k in t for k in ["robot", "로봇", "agv", "amr", "무인 운송", "용접"]):
        return "Robot/AMR"
    if any(k in t for k in ["vision", "비전", "검사", "품질"]):
        return "Vision/Inspection"
    if any(k in t for k in ["energy", "esg", "탄소", "전력", "fems"]):
        return "Energy/FEMS/ESG"
    if any(k in t for k in ["safety", "안전"]):
        return "Safety"
    if any(k in t for k in ["security", "보안", "mfa", "zone", "conduit"]):
        return "OT Security"
    if any(k in t for k in ["test", "measurement", "metrology", "계측", "측정"]):
        return "Test/Metrology"
    if any(k in t for k in ["지시", "workflow", "작업지시", "근태", "일원화"]):
        return "Ops/Workflow"
    if any(k in t for k in ["erp", "연동", "통합"]):
        return "Integration(ERP)"
    if any(k in t for k in ["%", "향상", "단축", "절감", "증가", "생산성", "공기"]):
        return "KPI/Outcome"
    return "Cloud/Data"


def _slice(md: str, start_hdr: str, end_hdr: str) -> str:
    s = md.find(start_hdr)
    if s == -1:
        return ''
    e = md.find(end_hdr, s + len(start_hdr))
    return md[s:e] if e != -1 else md[s:]


def parse_items(md: str) -> list[FactItem]:
    """Parse Fact/Inference bullets from sections we care about.

    We intentionally keep it heuristic and robust to formatting drift.
    """
    sections = [
        _slice(md, '## 2) 팩트 리스트', '## 3)'),
        _slice(md, '## 3) 생태계 맵 A', '## 4)'),
        _slice(md, '## 4) 생태계 맵 B', '## 5)'),
    ]

    items: list[FactItem] = []
    for sec in sections:
        for ln in sec.splitlines():
            t = ln.strip()
            if t.startswith('- Fact:'):
                items.append(FactItem(text=t[len('- Fact:'):].strip(), kind='fact'))
            elif t.startswith('- 추정'):
                items.append(FactItem(text=t[2:].strip(), kind='infer'))
    return items


def extract_kpi(md: str) -> dict[str, str]:
    # Heuristic for HD현대:
    # P: 생산성/리드타임/공기/생산능력
    # Q: 불량/품질
    # A: OEE/다운타임/예지
    # E: 에너지/탄소
    k = {'P': 'TBD', 'Q': 'TBD', 'A': 'TBD', 'E': 'TBD'}

    # P: pick lines with 생산성/단축/리드타임/생산능력 and %
    m = re.findall(r'[^\n]*?(생산성|노동\s*생산성|공기|기간|단축|리드타임|생산능력)[^\n]*?(\d+%|\d+→\d+)[^\n]*', md)
    if m:
        # Use a compact hand-written summary from known HD sources if present
        # Try to find the explicit known trio.
        if '노동 생산성' in md and '35%' in md and '56%' in md:
            k['P'] = '울산(건설기계): 노동생산성 +20%, 리드타임 -35%, 생산능력 +56%'
        elif '생산성 30%' in md and '공기 30%' in md:
            k['P'] = '조선(FOS 목표): 생산성 +30%, 공기 -30%'

    # Q: look for 품질 + 수치 else presence
    if '품질' in md or '검사' in md:
        k['Q'] = '비전/AI 품질검사 도입 (수치 TBD)'

    # A: look for 고장 예측/예지
    if '고장' in md or '예지' in md:
        k['A'] = '로봇 상태 감시/고장 예측(HI FACTORY) (수치 TBD)'

    # E: energy/esg
    if any(w in md.lower() for w in ['esg', '에너지', '탄소', 'fems']):
        k['E'] = 'ESG/에너지 관련 사례 존재 (수치 TBD)'

    return k


def add_sources_footer(slide, prs, left, urls: list[str]):
    ft = slide.shapes.add_textbox(left, prs.slide_height - Inches(0.4), prs.slide_width - left * 2, Inches(0.3))
    tf = ft.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED

    # prefix
    r0 = p.add_run()
    r0.text = 'Sources: '
    r0.font.size = Pt(9)
    r0.font.color.rgb = MUTED

    if not urls:
        r = p.add_run(); r.text = '-'
        r.font.size = Pt(9); r.font.color.rgb = MUTED
        return

    for i, url in enumerate(urls[:10]):
        disp = url.replace('https://', '').replace('http://', '')
        if disp.endswith('/'):
            disp = disp[:-1]
        run = p.add_run()
        run.text = ((' | ' if i else '') + disp)
        run.font.size = Pt(9)
        run.font.color.rgb = MUTED
        # Make each run clickable
        run.hyperlink.address = url


def add_detail_slide(prs, group: str, category: str, entries: list[FactItem], sources: list[str], back_slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_topbar(s, prs, f"{group} — {category} (detail)")

    margin = Inches(0.7)
    top = Inches(1.0)
    box_h = prs.slide_height - Inches(1.55)
    rect(s, margin, top, prs.slide_width - margin * 2, box_h, LIGHT, line_rgb=PRIMARY)

    # bullets
    tb = s.shapes.add_textbox(margin + Inches(0.3), top + Inches(0.25), prs.slide_width - margin * 2 - Inches(0.6), box_h - Inches(0.6))
    tf = tb.text_frame
    tf.clear(); tf.word_wrap = True

    if not entries:
        p = tf.paragraphs[0]
        p.text = '—'
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
    else:
        for i, it in enumerate(entries[:10]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            prefix = '(Fact) ' if it.kind == 'fact' else '(Inference) '
            p.text = prefix + it.text
            p.font.size = Pt(14 if i < 2 else 12)
            p.font.color.rgb = TEXT if it.kind == 'fact' else WARN

    # back button
    back = s.shapes.add_textbox(Inches(11.9), Inches(0.15), Inches(1.2), Inches(0.4))
    btf = back.text_frame; btf.clear()
    bp = btf.paragraphs[0]
    br = bp.add_run(); br.text = '↩ BACK'
    br.font.size = Pt(12); br.font.bold = True; br.font.color.rgb = WHITE
    back.fill.solid(); back.fill.fore_color.rgb = PRIMARY_DARK
    back.line.fill.background()
    back.click_action.target_slide = back_slide

    add_sources_footer(s, prs, margin, sources)
    return s


def add_matrix_slide(prs, group: str, items: list[FactItem], sources: list[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_topbar(s, prs, f"{group} — Dense Matrix (summary → click for detail)")

    # bucket items
    buckets = {c: [] for c in CATEGORIES}
    for it in items:
        cat = categorize(it.text)
        buckets.setdefault(cat, []).append(it)

    # Pre-create detail slides for non-empty categories
    detail_slides = {}
    for cat, entries in buckets.items():
        if entries:
            detail_slides[cat] = add_detail_slide(prs, group, cat, entries, sources, back_slide=s)

    # layout: 4x4 grid
    margin_x = Inches(0.5)
    margin_y = Inches(0.85)
    gap = Inches(0.15)
    grid_w = prs.slide_width - margin_x * 2
    grid_h = prs.slide_height - margin_y - Inches(0.55)
    cell_w = (grid_w - gap * 3) / 4
    cell_h = (grid_h - gap * 3) / 4

    cats = CATEGORIES[:16]
    for idx, cat in enumerate(cats):
        r = idx // 4
        c = idx % 4
        left = margin_x + c * (cell_w + gap)
        top = margin_y + r * (cell_h + gap)

        cell = rect(s, left, top, cell_w, cell_h, LIGHT, line_rgb=PRIMARY)
        if cat in detail_slides:
            cell.click_action.target_slide = detail_slides[cat]

        # title
        tb = s.shapes.add_textbox(left + Inches(0.12), top + Inches(0.05), cell_w - Inches(0.24), Inches(0.28))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = cat + (' ↗' if cat in detail_slides else '')
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = PRIMARY_DARK

        # summary body: 1-line best item
        body = s.shapes.add_textbox(left + Inches(0.12), top + Inches(0.33), cell_w - Inches(0.24), cell_h - Inches(0.40))
        tf = body.text_frame
        tf.clear(); tf.word_wrap = True

        entries = buckets.get(cat, [])
        if not entries:
            p = tf.paragraphs[0]
            p.text = '—'
            p.font.size = Pt(10)
            p.font.color.rgb = MUTED
        else:
            # prefer fact
            best = next((e for e in entries if e.kind == 'fact'), entries[0])
            p = tf.paragraphs[0]
            p.text = ('(F) ' if best.kind == 'fact' else '(I) ') + best.text
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT if best.kind == 'fact' else WARN

    add_sources_footer(s, prs, margin_x, sources)
    return s


def add_kpi_slide(prs, group: str, kpi: dict[str, str], sources: list[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_topbar(s, prs, f"{group} — KPI Snapshot (P/Q/A/E)")

    margin = Inches(0.7)
    top = Inches(1.05)
    w = prs.slide_width - margin * 2
    h = Inches(5.6)
    rect(s, margin, top, w, h, LIGHT, line_rgb=PRIMARY)

    # four rows
    row_h = h / 4
    lanes = [
        ('P', 'Productivity/Leadtime', kpi.get('P', 'TBD')),
        ('Q', 'Quality (Defect/Rework)', kpi.get('Q', 'TBD')),
        ('A', 'Asset (OEE/Downtime/PHM)', kpi.get('A', 'TBD')),
        ('E', 'Energy/ESG', kpi.get('E', 'TBD')),
    ]

    for i, (code, title, val) in enumerate(lanes):
        y = top + row_h * i
        # lane header box
        rect(s, margin, y, Inches(1.0), row_h, PRIMARY, line_rgb=PRIMARY)
        tb = s.shapes.add_textbox(margin + Inches(0.15), y + Inches(0.08), Inches(0.8), row_h - Inches(0.16))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = code
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = WHITE

        # title + value
        tb2 = s.shapes.add_textbox(margin + Inches(1.15), y + Inches(0.08), w - Inches(1.25), row_h - Inches(0.16))
        tf2 = tb2.text_frame
        tf2.clear()
        p1 = tf2.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_DARK
        p2 = tf2.add_paragraph()
        p2.text = val
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT

    ft = s.shapes.add_textbox(margin, prs.slide_height - Inches(0.4), prs.slide_width - margin * 2, Inches(0.3))
    tf = ft.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = 'Sources: ' + (' | '.join(sources) if sources else '-')
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def extract_sources(md: str, limit: int = 10) -> list[str]:
    """Return full URLs (kept clickable in PPT)."""
    src = []
    for ln in md.splitlines():
        t = ln.strip()
        if t.startswith('- Source:'):
            url = t.split(':', 1)[1].strip()
            if url:
                src.append(url)
    out = []
    seen = set()
    for u in src:
        if u in seen:
            continue
        seen.add(u)
        out.append(u)
        if len(out) >= limit:
            break
    return out


def add_framework_slide(prs, sources: list[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_topbar(s, prs, 'WEF 등대공장 — 기준/목표/확산 프레임 (내부 브리핑)')

    margin = Inches(0.7)
    top = Inches(1.0)
    w = prs.slide_width - margin * 2
    h = Inches(5.7)

    # 3 columns
    gap = Inches(0.25)
    col_w = (w - gap * 2) / 3

    boxes = [
        ('선정 기준(4요소)', [
            '기술(4IR Tech) 적용',
            '성과(Impact) — KPI 개선',
            '스케일(Replicability) — 복제 가능',
            '운영체계(사람/프로세스/데이터 거버넌스)',
        ]),
        ('KPI 4대 축(P/Q/A/E)', [
            'P 생산성/리드타임',
            'Q 품질(불량/재작업)',
            'A 설비(OEE/다운타임/예지)',
            'E 에너지/ESG',
        ]),
        ('확산 레이어(How)', [
            'Use-case 복제',
            'Platform 확산(MES+데이터)',
            'Operating model 확산(런북/RACI)',
            'Ecosystem 확산(협력사/SME)',
        ]),
    ]

    for i, (title, bullets) in enumerate(boxes):
        left = margin + i * (col_w + gap)
        rect(s, left, top, col_w, h, LIGHT, line_rgb=PRIMARY)
        tb = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), col_w - Inches(0.5), Inches(0.5))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK

        body = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.8), col_w - Inches(0.5), h - Inches(1.0))
        tf = body.text_frame; tf.clear(); tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = '• ' + b
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT

    add_sources_footer(s, prs, margin, sources)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, 'Smart Factory Solution Map — DENSE', '8개 그룹 | 16칸 요약 → 클릭하면 상세 슬라이드 | KPI(P/Q/A/E)')

    base = Path(__file__).resolve().parents[2] / 'sources'

    # Framework slide sources (use a few authoritative pages we already have access to)
    framework_sources = [
        'https://newsroom.posco.com/en/posco-named-lighthouse-factory-by-world-economic-forum/',
        'https://live.lge.co.kr/lg-wef-lighthousefactory/',
        'https://www.lsholdings.com/ko/media/news/385a523443424b4f5556485338745a73366346453871474d547866346b57776f',
    ]
    add_framework_slide(prs, framework_sources)

    groups = [
        ('HD현대', 'HD현대.md'),
        ('두산', '두산그룹.md'),
        ('현대차그룹', '현대자동차그룹.md'),
        ('삼성', '삼성.md'),
        ('LG', 'LG.md'),
        ('SK', 'SK.md'),
        ('포스코', '포스코.md'),
        ('LS그룹', 'LS그룹.md'),
    ]

    for g, fn in groups:
        md = (base / fn).read_text(encoding='utf-8', errors='ignore')
        items = parse_items(md)
        sources = extract_sources(md)
        kpi = extract_kpi(md)
        add_matrix_slide(prs, g, items, sources)
        add_kpi_slide(prs, g, kpi, sources)

    out_dir = Path(__file__).resolve().parents[2] / 'output'
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / 'smart-factory-solution-map_DENSE_blue.pptx'
    prs.save(out_path)
    print(out_path)


if __name__ == '__main__':
    main()

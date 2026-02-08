#!/usr/bin/env python3
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(hexstr: str) -> RGBColor:
    hexstr = hexstr.lstrip('#')
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

PRIMARY = rgb('#1E5EFF')
PRIMARY_DARK = rgb('#0B1B4D')
LIGHT = rgb('#EAF2FF')
TEXT = rgb('#111827')
WHITE = rgb('#FFFFFF')

SLIDE_RE = re.compile(r"^##\s+슬라이드\s+(\d+)\.\s*(.+?)\s*$")
NOTES_RE = re.compile(r"^발표자 노트:\s*$")


def parse_md(md_text: str):
    lines = md_text.splitlines()
    slides = []
    i = 0
    cur = None
    in_notes = False

    def finish():
        nonlocal cur
        if cur:
            # trim trailing empties
            cur['bullets'] = [b for b in cur['bullets'] if b.strip()]
            cur['notes'] = cur['notes'].strip()
            slides.append(cur)
            cur = None

    while i < len(lines):
        line = lines[i]
        m = SLIDE_RE.match(line)
        if m:
            finish()
            cur = {
                'num': int(m.group(1)),
                'title': m.group(2).strip(),
                'bullets': [],
                'notes': ''
            }
            in_notes = False
            i += 1
            continue

        if cur is None:
            i += 1
            continue

        if NOTES_RE.match(line):
            in_notes = True
            i += 1
            continue

        if in_notes:
            cur['notes'] += line + "\n"
        else:
            # bullets are '- ' lines; keep simple
            if line.strip().startswith('- '):
                cur['bullets'].append(line.strip()[2:].strip())
        i += 1

    finish()
    slides.sort(key=lambda s: s['num'])
    return slides


def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, prs, title: str):
    # blue bar across top
    bar_h = Inches(0.7)
    bar = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE (avoid importing enum for speed)
        0, 0, prs.slide_width, bar_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), prs.slide_width - Inches(1.2), Inches(0.5))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, prs, bullets):
    left = Inches(0.9)
    top = Inches(1.2)
    width = prs.slide_width - Inches(1.8)
    height = prs.slide_height - Inches(1.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)


def add_notes(slide, notes: str):
    if not notes.strip():
        return
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()
    tf.text = notes.strip()


def build_ppt(slides, out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide from slide 1
    for s in slides:
        if s['num'] == 1:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            set_slide_bg(slide, PRIMARY_DARK)
            # subtle light band
            band = slide.shapes.add_shape(1, 0, Inches(5.8), prs.slide_width, Inches(1.7))
            band.fill.solid(); band.fill.fore_color.rgb = PRIMARY
            band.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), prs.slide_width - Inches(1.8), Inches(2.0))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = s['bullets'][0] if s['bullets'] else s['title']
            r.font.size = Pt(44)
            r.font.bold = True
            r.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT

            # subtitle lines
            sub_lines = []
            if s['bullets']:
                sub_lines = s['bullets'][1:]
            if sub_lines:
                sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.5), prs.slide_width - Inches(1.8), Inches(2.0))
                stf = sub_box.text_frame
                stf.clear()
                for i, line in enumerate(sub_lines):
                    p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(24)
                    p.font.color.rgb = LIGHT
            add_notes(slide, s['notes'])
            break

    # Other slides
    for s in slides:
        if s['num'] == 1:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, WHITE)
        add_top_bar(slide, prs, f"슬라이드 {s['num']}. {s['title']}")
        add_bullets(slide, prs, s['bullets'])
        add_notes(slide, s['notes'])

        # small footer accent
        footer = slide.shapes.add_shape(1, 0, prs.slide_height - Inches(0.18), prs.slide_width, Inches(0.18))
        footer.fill.solid(); footer.fill.fore_color.rgb = LIGHT
        footer.line.fill.background()

    prs.save(out_path)


def main():
    md_path = Path(__file__).with_name('PPT_SLIDES.md')
    md = md_path.read_text(encoding='utf-8')
    slides = parse_md(md)
    out = Path(__file__).with_name('제조_AI_DX_용어집_웹앱_매뉴얼_블루테마.pptx')
    build_ppt(slides, out)
    print(str(out))


if __name__ == '__main__':
    main()
